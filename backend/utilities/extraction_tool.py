import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO

class ExtractionTool:
    def __init__(self):
        pass

    # Function to extract text from PDF files using fitz
    def extract_text_from_pdf(self, pdf_bytes: bytes) -> str:
        # Use the byte stream directly with PyMuPDF
        with fitz.open("pdf", pdf_bytes) as pdf_doc:
            text = ""
            for page_num in range(len(pdf_doc)):
                page = pdf_doc.load_page(page_num)
                text += page.get_text()
            return text

    # Function to extract text from PowerPoint files using Presentation
    def extract_text_from_pptx(self, file_bytes: bytes) -> str:
        text = ""
        # Use BytesIO to read the byte content as a file-like object
        presentation = Presentation(BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        return text
